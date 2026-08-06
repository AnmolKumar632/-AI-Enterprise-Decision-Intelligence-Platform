import os
import datetime
from utilities.custom_logger import get_logger

# ReportLab PDF imports
from reportlab.lib.pagesizes import letter
from reportlab.lib import colors
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle

# PPTX Presentation imports
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

logger = get_logger('report_generator_engine')

class ReportGenerator:
    def __init__(self, output_dir: str):
        self.output_dir = output_dir
        os.makedirs(self.output_dir, exist_ok=True)
        
    def generate_pdf(self, project_name: str, dataset_name: str, quality_score: float, 
                     model_name: str, model_metrics: dict, forecast_summary: str, 
                     anomaly_summary: str, recommendations: list) -> str:
        """Generate a professional executive summary PDF report."""
        pdf_filename = f"report_{project_name.lower().replace(' ', '_')}_{int(datetime.datetime.utcnow().timestamp())}.pdf"
        file_path = os.path.join(self.output_dir, pdf_filename)
        
        # Create SimpleDocTemplate
        doc = SimpleDocTemplate(
            file_path,
            pagesize=letter,
            rightMargin=54, leftMargin=54, topMargin=54, bottomMargin=54
        )
        
        styles = getSampleStyleSheet()
        
        # Custom styles for corporate presentation
        title_style = ParagraphStyle(
            'DocTitle',
            parent=styles['Normal'],
            fontName='Helvetica-Bold',
            fontSize=26,
            leading=30,
            textColor=colors.HexColor('#0F172A'), # Charcoal / Dark Slate
            spaceAfter=15
        )
        
        h1_style = ParagraphStyle(
            'SectionH1',
            parent=styles['Normal'],
            fontName='Helvetica-Bold',
            fontSize=16,
            leading=20,
            textColor=colors.HexColor('#1E3A8A'), # Navy Blue
            spaceBefore=15,
            spaceAfter=10,
            keepWithNext=True
        )
        
        body_style = ParagraphStyle(
            'BodyText',
            parent=styles['Normal'],
            fontName='Helvetica',
            fontSize=10,
            leading=14,
            textColor=colors.HexColor('#334155'), # Medium gray
            spaceAfter=8
        )
        
        bold_body_style = ParagraphStyle(
            'BoldBodyText',
            parent=body_style,
            fontName='Helvetica-Bold'
        )
        
        story = []
        
        # 1. Header / Cover Block
        story.append(Paragraph("AEDIP Executive Decision Intelligence Report", title_style))
        story.append(Paragraph(f"<b>Project Workspace:</b> {project_name}", body_style))
        story.append(Paragraph(f"<b>Generated On:</b> {datetime.date.today().strftime('%B %d, %Y')}", body_style))
        story.append(Spacer(1, 20))
        
        # 2. Executive Summary
        story.append(Paragraph("1. Executive Summary", h1_style))
        exec_text = (
            "This decision intelligence report was automatically compiled by the AI Enterprise Decision "
            "Intelligence Platform (AEDIP). It synthesizes automated data cleaning profiles, machine learning model "
            "rankings, future-state time series forecasts, and unsupervised anomaly audits to provide stakeholders "
            "with empirical evidence to guide strategic directives."
        )
        story.append(Paragraph(exec_text, body_style))
        story.append(Spacer(1, 10))
        
        # 3. Dataset Profile & Quality Summary
        story.append(Paragraph("2. Dataset Hygiene & Quality Assessment", h1_style))
        story.append(Paragraph(f"<b>Dataset Analyzed:</b> {dataset_name}", body_style))
        story.append(Paragraph(f"<b>Data Hygiene Score:</b> {quality_score}/100", body_style))
        
        # Embed quality badge indicators
        quality_status = "Excellent" if quality_score >= 85 else "Fair" if quality_score >= 60 else "Critical Action Required"
        story.append(Paragraph(f"<b>Hygiene Classification:</b> {quality_status}", body_style))
        story.append(Spacer(1, 10))
        
        # 4. Model Leaderboard Metrics
        story.append(Paragraph("3. Predictive Modeling Performance", h1_style))
        story.append(Paragraph(f"<b>Optimal Algorithm Selected:</b> {model_name}", body_style))
        
        # Build metrics table
        metrics_data = [["Evaluation Metric", "Metric Value"]]
        for k, v in model_metrics.items():
            if isinstance(v, (int, float)):
                metrics_data.append([k.replace('_', ' ').capitalize(), str(round(v, 4))])
                
        t = Table(metrics_data, colWidths=[200, 150])
        t.setStyle(TableStyle([
            ('BACKGROUND', (0,0), (-1,0), colors.HexColor('#1E3A8A')),
            ('TEXTCOLOR', (0,0), (-1,0), colors.whitesmoke),
            ('ALIGN', (0,0), (-1,-1), 'LEFT'),
            ('FONTNAME', (0,0), (-1,0), 'Helvetica-Bold'),
            ('BOTTOMPADDING', (0,0), (-1,0), 6),
            ('BACKGROUND', (0,1), (-1,-1), colors.HexColor('#F8FAFC')),
            ('GRID', (0,0), (-1,-1), 1, colors.HexColor('#CBD5E1')),
            ('FONTNAME', (0,1), (-1,-1), 'Helvetica'),
            ('FONTSIZE', (0,0), (-1,-1), 9),
        ]))
        story.append(t)
        story.append(Spacer(1, 15))
        
        # Page break before Advanced Diagnostics
        story.append(PageBreak())
        
        # 5. Forecasting Insights
        story.append(Paragraph("4. Time-Series Forecasting Trend", h1_style))
        story.append(Paragraph(forecast_summary, body_style))
        story.append(Spacer(1, 10))
        
        # 6. Anomaly/Risk Profile
        story.append(Paragraph("5. Unsupervised Risk & Anomaly Assessment", h1_style))
        story.append(Paragraph(anomaly_summary, body_style))
        story.append(Spacer(1, 10))
        
        # 7. Strategic Recommendations
        story.append(Paragraph("6. Corporate Strategic Directives", h1_style))
        for rec in recommendations:
            story.append(Paragraph(f"• {rec}", body_style))
            
        # Build Document
        doc.build(story)
        logger.info(f"PDF report generated successfully at: {file_path}")
        return file_path
        
    def generate_pptx(self, project_name: str, dataset_name: str, quality_score: float,
                      model_name: str, model_metrics: dict, recommendations: list) -> str:
        """Generate a corporate PowerPoint slideshow (.pptx)."""
        pptx_filename = f"presentation_{project_name.lower().replace(' ', '_')}_{int(datetime.datetime.utcnow().timestamp())}.pptx"
        file_path = os.path.join(self.output_dir, pptx_filename)
        
        prs = Presentation()
        
        # Custom color definitions
        navy_color = RGBColor(30, 58, 138)
        charcoal_color = RGBColor(15, 23, 42)
        gray_color = RGBColor(100, 116, 139)
        
        # 1. COVER SLIDE
        slide_layout = prs.slide_layouts[0] # Title Layout
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "AI Enterprise Decision Analytics"
        subtitle.text = f"Workspace: {project_name}\nCompiled: {datetime.date.today().strftime('%d %B %Y')}"
        
        title.text_frame.paragraphs[0].font.color.rgb = navy_color
        title.text_frame.paragraphs[0].font.bold = True
        
        # 2. DATASET SLIDE
        slide_layout = prs.slide_layouts[1] # Content Layout
        slide = prs.slides.add_slide(slide_layout)
        shapes = slide.shapes
        shapes.title.text = "Dataset Profile & Hygiene Assessment"
        shapes.title.text_frame.paragraphs[0].font.color.rgb = navy_color
        
        tf = shapes.placeholders[1].text_frame
        tf.text = f"Analyzing dataset: {dataset_name}"
        
        p = tf.add_paragraph()
        p.text = f"Data Quality Score: {quality_score}/100"
        p.font.bold = True
        p.font.size = Pt(20)
        p.font.color.rgb = navy_color if quality_score >= 75 else RGBColor(185, 28, 28)
        
        p2 = tf.add_paragraph()
        p2.text = (
            "• Outlier clipping and statistical imputations applied to stabilize covariance matrices.\n"
            "• Schema features scanned for automated predictive modeling ingestion.\n"
            "• Missing values resolved using median and modal imputation vectors."
        )
        p2.font.size = Pt(16)
        p2.font.color.rgb = charcoal_color
        
        # 3. PREDICTIVE MODELLING SLIDE
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = "AutoML Leaderboard Results"
        slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = navy_color
        
        tf = slide.shapes.placeholders[1].text_frame
        tf.text = f"Optimal Algorithm selected: {model_name}"
        
        p = tf.add_paragraph()
        p.text = "Top Model Validation Metrics:"
        p.font.bold = True
        p.font.size = Pt(16)
        
        for k, v in model_metrics.items():
            if isinstance(v, (int, float)):
                p_m = tf.add_paragraph()
                p_m.text = f"  - {k.replace('_', ' ').capitalize()}: {round(v, 4)}"
                p_m.font.size = Pt(14)
                p_m.font.color.rgb = charcoal_color
                
        # 4. RECOMMENDATIONS SLIDE
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = "Strategic Executive Directives"
        slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = navy_color
        
        tf = slide.shapes.placeholders[1].text_frame
        tf.text = "Corporate action items derived from analytical findings:"
        tf.text_frame.paragraphs[0].font.bold = True
        
        for rec in recommendations:
            p_r = tf.add_paragraph()
            p_r.text = f"• {rec}"
            p_r.font.size = Pt(14)
            p_r.font.color.rgb = charcoal_color
            
        prs.save(file_path)
        logger.info(f"PowerPoint report generated successfully at: {file_path}")
        return file_path
